import os
import sys
import subprocess

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Error: 'python-pptx' package not installed. Installing automatically via pip...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Configure widescreen 16:9 dimensions (13.333 inches x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors (Geospatial Command Center Dark Theme)
    BG_COLOR = RGBColor(11, 19, 43)        # #0B132B (Deep Obsidian / Night Sky)
    PANEL_COLOR = RGBColor(28, 37, 65)     # #1C2541 (Slate Blue Panel)
    TEXT_WHITE = RGBColor(255, 255, 255)   # #FFFFFF
    TEXT_SILVER = RGBColor(220, 230, 242)  # #E0E6ED
    CYAN_ACCENT = RGBColor(0, 245, 212)    # #00F5D4 (GeoAI Intelligence / Water)
    AMBER_ACCENT = RGBColor(255, 190, 11)  # #FFBE0B (Warning / Moderate)
    RED_ACCENT = RGBColor(255, 76, 76)     # #FF4C4C (Critical Extreme Alert)

    blank_layout = prs.slide_layouts[6] # completely clean blank layout

    def apply_slide_base(slide, title_text):
        # Background rectangle
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # Top Header Cyan Line
        header_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.85), Inches(13.333), Inches(0.05))
        header_line.fill.solid()
        header_line.fill.fore_color.rgb = CYAN_ACCENT
        header_line.line.fill.background()

        # Title Text box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.65))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text.upper()
        p.font.name = 'Segoe UI'
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE

        # Footer Text box
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.35))
        ftf = footer.text_frame
        fp = ftf.paragraphs[0]
        fp.text = "GEONARRATIVE AI  |  M.SC. THESIS DEFENSE  |  SYMBIOSIS INSTITUTE OF GEOINFORMATICS (SIG), PUNE"
        fp.font.name = 'Segoe UI'
        fp.font.size = Pt(10)
        fp.font.color.rgb = TEXT_SILVER

    def add_speaker_notes(slide, notes_text):
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes_text

    print("Synthesizing Slide 1: Title Gateway...")
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = BG_COLOR
    bg1.line.fill.background()
    
    # Title decorative accent frame
    accent1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(1.8), Inches(0.15), Inches(3.2))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = CYAN_ACCENT
    accent1.line.fill.background()

    tbox1 = slide1.shapes.add_textbox(Inches(1.4), Inches(1.7), Inches(11), Inches(4.5))
    tf1 = tbox1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "GEONARRATIVE AI"
    p1.font.name = 'Segoe UI'
    p1.font.size = Pt(48)
    p1.font.bold = True
    p1.font.color.rgb = CYAN_ACCENT
    
    p2 = tf1.add_paragraph()
    p2.text = "An Intelligent Geospatial Digital Twin and GeoAI Decision-Support Framework for Urban Flood Risk Assessment in Pune"
    p2.font.name = 'Segoe UI'
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_WHITE
    p2.space_before = Pt(15)

    p3 = tf1.add_paragraph()
    p3.text = "\nStudent: Harnish Savant\nProgramme: M.Sc. Data Science and Spatial Analytics (Geo-Intelligence)\nInstitute: Symbiosis Institute of Geoinformatics (SIG), Pune"
    p3.font.name = 'Segoe UI'
    p3.font.size = Pt(16)
    p3.font.color.rgb = TEXT_SILVER
    p3.space_before = Pt(25)
    
    add_speaker_notes(slide1, "Good day, members of the evaluation committee and esteemed faculty. Today, I am proud to present my M.Sc. dissertation research titled GeoNarrative AI. This research addresses one of Pune's most pressing urban challenges—flooding—by bridging rigorous spatial analytical modeling with an intelligent, 3D WebGL Digital Twin and an explainable GeoAI decision-support platform.")

    print("Synthesizing Slide 2: Problem Definition...")
    slide2 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide2, "1. Problem Definition: The Urban Hydrological Crisis in Pune")
    box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "Accelerated Urbanization & Monsoon Cloudburst Vulnerability"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = CYAN_ACCENT
    
    bullets2 = [
        ("Topographical Complexity", "The Mula-Mutha river valley experiences intense seasonal monsoon cloudbursts. Extreme urbanization and increased impervious surface sealing lead to rapid runoff accumulation in low-lying basins."),
        ("Limitations of Traditional Static Maps", "Standard 2D GIS flood hazard printouts lack real-time dynamic interactivity and fail to provide municipal responders with exploratory scenario simulation."),
        ("Operational Data Silos", "Hydrological spatial modeling (AHP/GIS) remains completely decoupled from operational municipal emergency disaster coordination engines."),
        ("Lack of Analytical Explainability", "Emergency coordinators lack automated, explainable decision-support interfaces to interpret spatial numerical susceptibility indices instantaneously during critical flood crises.")
    ]
    for head, body in bullets2:
        p_sub = tf2.add_paragraph()
        p_sub.text = f"•  {head}: "
        p_sub.font.size = Pt(18); p_sub.font.bold = True; p_sub.font.color.rgb = AMBER_ACCENT
        p_sub.space_before = Pt(12)
        run = p_sub.add_run()
        run.text = body
        run.font.size = Pt(17); run.font.bold = False; run.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide2, "As Pune experiences unprecedented urban expansion, the proliferation of impervious surfaces alongside erratic monsoon precipitation has turned urban flood management into a critical priority. Traditional disaster management relies heavily on static, two-dimensional GIS printouts or decoupled tables. These conventional deliverables fail to provide emergency controllers with interactive multi-scenario exploration or real-time explainability during critical flood events.")

    print("Synthesizing Slide 3: Research Objectives (RQ1-RQ3)...")
    slide3 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide3, "2. Research Objectives (RQ1 - RQ3) & Scientific Novelty")
    box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.4))
    tf3 = box3.text_frame
    tf3.word_wrap = True
    
    rq_data = [
        ("RQ1: Multi-Criteria Spatial Hazard Modeling", "Formulate and mathematically validate a continuous Flood Susceptibility Index (FSI) across Pune utilizing a robust 5-factor Analytic Hierarchy Process (AHP) weighted linear overlay.", CYAN_ACCENT),
        ("RQ2: Infrastructure Exposure Quantification", "Establish a comprehensive, multi-scenario evaluation matrix evaluating vulnerability and physical exposure across municipal structural buildings and transportation networks.", AMBER_ACCENT),
        ("RQ3: Explainable Digital Twin Synthesis", "Develop a real-time, interactive 3D WebGL Digital Twin integrated with an LLM-driven GeoAI natural language assistant to automate technical disaster situation reports.", RGBColor(120, 160, 255))
    ]
    for i, (title, desc, color) in enumerate(rq_data):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = title
        p.font.size = Pt(21); p.font.bold = True; p.font.color.rgb = color
        p.space_before = Pt(20) if i > 0 else Pt(0)
        p_desc = tf3.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(17); p_desc.font.color.rgb = TEXT_WHITE
        p_desc.space_before = Pt(5)
    add_speaker_notes(slide3, "To resolve these operational bottlenecks, this thesis establishes three comprehensive research objectives. First, to model flood susceptibility across Pune using a rigorous five-factor Analytic Hierarchy Process. Second, to quantify actual infrastructure exposure across multiple storm severities. And third, to engineer an enterprise-grade 3D Digital Twin that integrates an LLM-driven GeoAI assistant capable of synthesizing automated disaster reports and explaining spatial analytics through natural language dialogue.")

    print("Synthesizing Slide 4: Study Area Scope...")
    slide4 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide4, "3. Study Area Scope: Pune Municipal Corporation (PMC)")
    box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(6.0), Inches(5.5))
    tf4 = box4.text_frame
    tf4.word_wrap = True
    p = tf4.paragraphs[0]
    p.text = "Empirical Infrastructure Scope & Inventory"
    p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = CYAN_ACCENT
    
    metrics = [
        ("Administrative Study Area", "506.91 sq. km", "Entire PMC domain situated on the Deccan Plateau."),
        ("Structural Building Inventory", "180,307 structures", "Georegistered structural building footprints used for exposure analysis."),
        ("Municipal Road Network", "7,445.90 km", "Multimodal transport and urban vehicular arterial roads.")
    ]
    for title, val, comment in metrics:
        p_m = tf4.add_paragraph()
        p_m.text = f"\n• {title}: "
        p_m.font.size = Pt(17); p_m.font.bold = True; p_m.font.color.rgb = TEXT_SILVER
        run1 = p_m.add_run()
        run1.text = val
        run1.font.size = Pt(18); run1.font.bold = True; run1.font.color.rgb = AMBER_ACCENT
        p_c = tf4.add_paragraph()
        p_c.text = f"   ({comment})"
        p_c.font.size = Pt(14); p_c.font.color.rgb = TEXT_WHITE

    # Embed Road Network & Building distribution if available
    img_road = os.path.join("Data", "png", "2.jpg")
    img_build = os.path.join("Data", "png", "4.jpg")
    if os.path.exists(img_road):
        slide4.shapes.add_picture(img_road, Inches(6.8), Inches(1.2), width=Inches(6.0))
    if os.path.exists(img_build):
        slide4.shapes.add_picture(img_build, Inches(6.8), Inches(4.1), width=Inches(6.0))
    add_speaker_notes(slide4, "Our domain of investigation covers the entire administrative boundary of the Pune Municipal Corporation, spanning exactly 506.91 square kilometers. To perform high-resolution exposure assessment, we integrated and pre-processed a massive municipal inventory comprising over 7,445 kilometers of road networks and exactly 180,307 georegistered building structures, creating a robust baseline for vulnerability evaluation.")

    print("Synthesizing Slide 5: Methodological Framework...")
    slide5 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide5, "4. Methodological Framework: UNDRR Disaster Risk Standard")
    box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf5 = box5.text_frame
    tf5.word_wrap = True
    p = tf5.paragraphs[0]
    p.text = "United Nations UNDRR Risk Formulation:"
    p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = CYAN_ACCENT
    
    p_eq = tf5.add_paragraph()
    p_eq.text = "Risk  =  Hazard  ×  Exposure  ×  Vulnerability"
    p_eq.font.size = Pt(28); p_eq.font.bold = True; p_eq.font.color.rgb = RED_ACCENT
    p_eq.space_before = Pt(15); p_eq.alignment = PP_ALIGN.CENTER
    
    steps = [
        ("1. Geospatial ETL & Data Harmonization", "Ingestion of OpenStreetMap spatial layers, USGS EarthEngine topography (DEM), and hydrological streamlines into PostGIS."),
        ("2. Analytical Raster Algebra (GDAL / AHP)", "Derivation of Euclidean river distance arrays, topographic slope angles, and 5-factor weighted overlay algebra."),
        ("3. Infrastructure Spatial Intersection", "Multi-scenario spatial intersection mapping building structural impact and road network severance."),
        ("4. Digital Twin WebGL & GeoAI Synthesis", "3D CesiumJS visual scene generation and Large Language Model (LLM) natural language analytical explanation.")
    ]
    for step, desc in steps:
        p_s = tf5.add_paragraph()
        p_s.text = step
        p_s.font.size = Pt(18); p_s.font.bold = True; p_s.font.color.rgb = AMBER_ACCENT
        p_s.space_before = Pt(15)
        p_d = tf5.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(16); p_d.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide5, "Our methodology is structurally grounded in the United Nations Disaster Risk Reduction framework, where overall disaster risk is defined as the mathematical product of Hazard intensity, Structural Exposure, and localized Vulnerability. By constructing a fully end-to-end analytical pipeline—from PostGIS database ingestion to GDAL spatial raster algebra—our platform guarantees that every visual layer displayed in the 3D frontend is underpinned by defensible statistical equations.")

    print("Synthesizing Slide 6: AHP Weighted Overlay Engine...")
    slide6 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide6, "5. Hazard Modeling: The 5-Factor AHP Overlay Engine")
    box6 = slide6.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.1), Inches(5.5))
    tf6 = box6.text_frame; tf6.word_wrap = True
    p6 = tf6.paragraphs[0]
    p6.text = "Mathematical Formulation:   FSI = ∑ ( w_j · S_ij )   [where w_j are normalized factor weights]"
    p6.font.size = Pt(20); p6.font.bold = True; p6.font.color.rgb = CYAN_ACCENT
    
    weights = [
        ("Digital Elevation Model (DEM)", "0.35 (35%)", "Primary hydrological baseline governing gravitational flow density and drainage aggregation."),
        ("Euclidean Distance to River", "0.25 (25%)", "Direct geographical distance to perennial Mula, Mutha, and Pavana overflow alignments."),
        ("Topographic Slope Angle", "0.20 (20%)", "Surface gradient influencing water velocity and localized ponding accumulation in basin flats."),
        ("Land Use / Land Cover (LULC)", "0.10 (10%)", "Impervious surface runoff coefficients and natural infiltration inhibition."),
        ("Structural Building Density", "0.10 (10%)", "Anthropogenic urban structural sealing and artificial surface flow channeling.")
    ]
    for w_title, w_val, w_desc in weights:
        pw = tf6.add_paragraph()
        pw.text = f"•  {w_title}  →  Weight: "
        pw.font.size = Pt(17); pw.font.bold = True; pw.font.color.rgb = TEXT_SILVER
        pw.space_before = Pt(12)
        run = pw.add_run()
        run.text = w_val
        run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = AMBER_ACCENT
        run2 = pw.add_run()
        run2.text = f"   ({w_desc})"
        run2.font.size = Pt(15); run2.font.bold = False; run2.font.color.rgb = TEXT_WHITE

    pval = tf6.add_paragraph()
    pval.text = "\nMathematical Consistency Verification:  Consistency Ratio (CR) = 0.038  (Complying strictly with Saaty's academic threshold of CR < 0.10)."
    pval.font.size = Pt(17); pval.font.bold = True; pval.font.color.rgb = RED_ACCENT
    add_speaker_notes(slide6, "To compute the continuous Flood Susceptibility Index, we employed an Analytic Hierarchy Process utilizing a five-factor weighted linear combination. Based on geomorphological dominance, Topographic Elevation and Euclidean Distance to permanent rivers received the highest weights of 35% and 25%, respectively, followed by Topographic Slope at 20%, and LULC alongside Building Density at 10% each. Our expert pairwise judgment matrix achieved a rigorous Consistency Ratio of just 0.038, well below the acceptable limit of 0.10, proving mathematical soundness.")

    print("Synthesizing Slide 7: Geospatial Factor Mapping...")
    slide7 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide7, "6. Geospatial Factor Mapping: Morphological & Hydrology Drivers")
    box7 = slide7.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.0), Inches(0.8))
    tf7 = box7.text_frame; tf7.word_wrap = True
    tf7.paragraphs[0].text = "Spatial Driver Outputs (Elevation, Slope & River Distance Buffering):"
    tf7.paragraphs[0].font.size = Pt(20); tf7.paragraphs[0].font.bold = True; tf7.paragraphs[0].font.color.rgb = CYAN_ACCENT
    
    img_dem = os.path.join("Data", "png", "6.jpg")
    img_slope = os.path.join("Data", "png", "8.jpg")
    img_river = os.path.join("Data", "png", "9.jpg")
    if os.path.exists(img_dem):
        slide7.shapes.add_picture(img_dem, Inches(0.4), Inches(1.9), width=Inches(4.0))
    if os.path.exists(img_slope):
        slide7.shapes.add_picture(img_slope, Inches(4.6), Inches(1.9), width=Inches(4.0))
    if os.path.exists(img_river):
        slide7.shapes.add_picture(img_river, Inches(8.8), Inches(1.9), width=Inches(4.0))
    
    t7 = slide7.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.2), Inches(1.0))
    p_cap = t7.text_frame.paragraphs[0]
    p_cap.text = "• DEM: Elevation ranges from <540m along river beds to >1,100m in western ghat hills.\n• Slope: Ultra-flat accumulation basins identified in Shivajinagar and Aundh.\n• Distance to River: Euclidean buffering demonstrates localized structural vulnerability zones along main channels."
    p_cap.font.size = Pt(15); p_cap.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide7, "Here we inspect the core spatial raster outputs generated during our GIS modeling phase. On the left, the Digital Elevation Model illustrates the steep elevation drop from the western hills into the Mula-Mutha valley. In the center, our Topographic Slope map clearly delineates ultra-low gradient floodplains where water ponding inevitably concentrates during prolonged storms. On the right, our continuous Euclidean Distance matrix establishes precise spatial proximity bands from every building to the primary river network.")

    print("Synthesizing Slide 8: FSI Results...")
    slide8 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide8, "7. Flood Susceptibility Index (FSI) & Structural Density")
    box8 = slide8.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(12.0), Inches(0.8))
    tf8 = box8.text_frame; tf8.word_wrap = True
    tf8.paragraphs[0].text = "Synthesized Flood Susceptibility vs. Anthropogenic Structural Building Density:"
    tf8.paragraphs[0].font.size = Pt(20); tf8.paragraphs[0].font.bold = True; tf8.paragraphs[0].font.color.rgb = CYAN_ACCENT

    img_fsi = os.path.join("Data", "png", "11.jpg")
    img_dense = os.path.join("Data", "png", "10.jpg")
    if os.path.exists(img_fsi):
        slide8.shapes.add_picture(img_fsi, Inches(0.5), Inches(1.8), width=Inches(5.9))
    if os.path.exists(img_dense):
        slide8.shapes.add_picture(img_dense, Inches(6.8), Inches(1.8), width=Inches(5.9))
        
    t8 = slide8.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12.2), Inches(1.0))
    p_cap8 = t8.text_frame.paragraphs[0]
    p_cap8.text = "• High/Very High Susceptibility (Red/Purple) concentrates heavily in riparian wards along Mula-Mutha confluence.\n• Overlaying structural building density (Right) reveals severe urban exposure in historic high-density civic cores."
    p_cap8.font.size = Pt(16); p_cap8.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide8, "By executing our AHP overlay algebra, we generate Figure 5.11—the definitive Flood Susceptibility Map of Pune Municipal Corporation. The spatial distribution demonstrates that high and very high susceptibility zones—shown in purple and red—cluster aggressively along river corridors and historic dense core wards. When superimposed against our building density map in Figure 5.10, we clearly visualize how extreme anthropogenic structural density exacerbates hydrological vulnerability.")

    print("Synthesizing Slide 9: Scientific Guardrails...")
    slide9 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide9, "8. Scientific Integrity: Inundation Depth Definition & Scope")
    box9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf9 = box9.text_frame; tf9.word_wrap = True
    p9 = tf9.paragraphs[0]
    p9.text = "CRITICAL METHODOLOGICAL TERMINOLOGY & GUARDRAILS"
    p9.font.size = Pt(22); p9.font.bold = True; p9.font.color.rgb = RED_ACCENT
    
    guard_bullets = [
        ("Distinction from Empirical Hydrodynamics", "The computational representations in GeoNarrative AI do NOT constitute direct numerical solutions to fluid Navier-Stokes or 2D Saint-Venant hydraulic shallow water physical equations (e.g., SWMM hydraulic modeling)."),
        ("Scenario-Derived Relative Inundation Depth Estimates", "All simulated flood depths represent 'scenario-derived relative inundation depth estimates' derived from topographically evaluated susceptibility indices, localized DEM depressions, river proximity decay coefficients, and precipitation scaling."),
        ("Operational Decision-Support Purpose", "This topological modeling architecture is explicitly designed for instantaneous, real-time spatial prioritization and early warning disaster coordination without the extreme computational CPU rendering latency associated with physical hydrodynamic simulations.")
    ]
    for head, body in guard_bullets:
        p_g = tf9.add_paragraph()
        p_g.text = f"\n•  {head}:"
        p_g.font.size = Pt(19); p_g.font.bold = True; p_g.font.color.rgb = CYAN_ACCENT
        p_g.space_before = Pt(10)
        p_b = tf9.add_paragraph()
        p_b.text = f"    {body}"
        p_b.font.size = Pt(17); p_b.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide9, "At this juncture, it is critical to state the scientific boundaries of our simulation framework with unconditional precision. The flood depth approximations generated within GeoNarrative AI represent scenario-derived relative inundation depth estimates. Rather than solving fluid mechanical Navier-Stokes or shallow water differential equations, our model derives relative inundation intensities from topological DEM depressions, susceptibility indices, and precipitation scaling. This provides municipal emergency commanders with instant, computationally efficient spatial prioritization without the extreme CPU rendering latency of traditional fluid modeling.")

    print("Synthesizing Slide 10: Multi-Scenario Exposure Centerpiece Table...")
    slide10 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide10, "9. Multi-Scenario Infrastructure Exposure Summary Matrix")
    box10 = slide10.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.0), Inches(0.6))
    box10.text_frame.word_wrap = True
    box10.text_frame.paragraphs[0].text = "Quantified Empirical Exposure across Normal, Moderate, Heavy, & Extreme Storms:"
    box10.text_frame.paragraphs[0].font.size = Pt(18); box10.text_frame.paragraphs[0].font.bold = True; box10.text_frame.paragraphs[0].font.color.rgb = CYAN_ACCENT

    # Add Table
    rows = 5; cols = 6
    table_shape = slide10.shapes.add_table(rows, cols, Inches(0.5), Inches(1.7), Inches(12.33), Inches(4.5))
    table = table_shape.table
    
    headers = ["Scenario", "Rainfall Rate", "Relative Inundation", "Inundated Area (sq km)", "Impacted Buildings", "Severed Road Length"]
    for i, head in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = head
        for p in cell.text_frame.paragraphs:
            p.font.name = 'Segoe UI'; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = CYAN_ACCENT
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = PANEL_COLOR

    row_data = [
        ("Normal", "0 mm/hr", "Baseline", "0.00 sq km (0.0%)", "0 (0.0%)", "0.00 km (0.0%)", TEXT_WHITE),
        ("Moderate", "30 mm/hr", "≤ 1.0 m (Ponding)", "~25.40 sq km (5.01%)", "~3,600 (2.00%)", "~370.00 km (4.97%)", TEXT_SILVER),
        ("Heavy", "65 mm/hr", "1.0 – 2.0 m (Waterlogging)", "~68.10 sq km (13.43%)", "~14,400 (7.99%)", "~930.00 km (12.49%)", AMBER_ACCENT),
        ("Extreme", "120 mm/hr", "> 2.0 m (Critical Flood)", "133.97 sq km (26.43%)", "40,723 (22.59%)", "1,877.47 km (25.21%)", RED_ACCENT)
    ]
    for row_idx, r in enumerate(row_data):
        for col_idx in range(6):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = str(r[col_idx])
            for p in cell.text_frame.paragraphs:
                p.font.name = 'Segoe UI'; p.font.size = Pt(15); p.font.bold = (row_idx == 3)
                p.font.color.rgb = r[6] if (col_idx == 0 or row_idx == 3) else TEXT_WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(18, 25, 50) if row_idx % 2 == 0 else PANEL_COLOR

    add_speaker_notes(slide10, "This comprehensive data matrix represents the central quantitative results of my dissertation. Across our four meteorological scenarios, we observe a dramatic scaling of infrastructure exposure. Under the Extreme Cloudburst Scenario—representing rainfall rates reaching 120 millimeters per hour—the model reveals that 133.97 square kilometers of urban area (over 26% of Pune) is inundated. In this critical regime, exactly 40,723 structural buildings (22.59% of the structural inventory) and 1,877.47 kilometers of municipal transport road networks (25.21% of total arterial roads) experience severe hazard exposure.")

    print("Synthesizing Slide 11: Enterprise Architecture & Digital Twin...")
    slide11 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide11, "10. Enterprise Software Architecture: 3D Digital Twin")
    box11 = slide11.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf11 = box11.text_frame; tf11.word_wrap = True
    tf11.paragraphs[0].text = "Full-Stack Geospatial Intelligence & WebGL Engine:"
    tf11.paragraphs[0].font.size = Pt(22); tf11.paragraphs[0].font.bold = True; tf11.paragraphs[0].font.color.rgb = CYAN_ACCENT
    
    twin_data = [
        ("Backend Analytical Pipeline (FastAPI & PostGIS)", "Asynchronous Python FastAPI microservices managing spatial intersecting queries over PostGIS municipal databases, streaming precomputed 45-frame raster animations at 60 FPS."),
        ("Frontend 3D Digital Twin (CesiumJS WebGL)", "High-fidelity terrain visualization blending realistic 3D building structural tiles with persistent, dynamic permanent river streamlines."),
        ("Dynamic Infrastructure Hazard Color-Ramps", "Real-time visual geometry shader injection shifting structural building styling dynamically from baseline Safe Cyan to Critical Hazard Red based on localized flood progression.")
    ]
    for t_head, t_body in twin_data:
        p_t = tf11.add_paragraph()
        p_t.text = f"\n•  {t_head}:"
        p_t.font.size = Pt(19); p_t.font.bold = True; p_t.font.color.rgb = AMBER_ACCENT
        p_b = tf11.add_paragraph()
        p_b.text = f"    {t_body}"
        p_b.font.size = Pt(17); p_b.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide11, "To operationalize these spatial findings, we engineered a state-of-the-art enterprise WebGL software architecture. Our backend combines Python FastAPI microservices with a PostGIS spatial database engine to execute high-performance spatial intersecting queries in real time. Our frontend leverages CesiumJS 3D WebGL rendering, incorporating pre-compiled animation frames to visualize flood propagation smoothly across the terrain while dynamically shifting structural building colors from safe cyan to critical red as simulated water exposure intensifies.")

    print("Synthesizing Slide 12: GeoAI Assistant...")
    slide12 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide12, "11. Explainable GeoAI Conversational Assistant & Report Agent")
    box12 = slide12.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf12 = box12.text_frame; tf12.word_wrap = True
    tf12.paragraphs[0].text = "LLM-Driven Natural Language Decision Support:"
    tf12.paragraphs[0].font.size = Pt(22); tf12.paragraphs[0].font.bold = True; tf12.paragraphs[0].font.color.rgb = CYAN_ACCENT

    ai_data = [
        ("Natural Language Spatial Querying", "Translates non-technical questions (e.g., 'Which ward has the most buildings at risk in heavy storms?') directly into PostGIS analytical queries and presents readable responses."),
        ("Automated Situation Report Agent", "Instantly compiles active simulated parameters into structured, downloadable emergency situation reports and executive ward-level evacuation briefs."),
        ("Bridging Academic GIS and Emergency Operations", "Removes technical barriers, enabling municipal disaster response commanders to operate advanced geospatial risk models without coding or GIS software expertise.")
    ]
    for a_head, a_body in ai_data:
        p_a = tf12.add_paragraph()
        p_a.text = f"\n•  {a_head}:"
        p_a.font.size = Pt(19); p_a.font.bold = True; p_a.font.color.rgb = RGBColor(120, 160, 255)
        p_ab = tf12.add_paragraph()
        p_ab.text = f"    {a_body}"
        p_ab.font.size = Pt(17); p_ab.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide12, "A defining novelty of GeoNarrative AI is the integration of an explainable Large Language Model conversational assistant. In a real-world municipal control room, disaster management coordinators rarely have GIS programming expertise. Our GeoAI assistant enables coordinators to simply type questions such as 'What is the structural impact on Shivajinagar during an extreme cloudburst?' to instantly retrieve verified empirical figures and generate formal, automated executive situation reports in seconds.")

    print("Synthesizing Slide 13: Technical Limitations & Future Scope...")
    slide13 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide13, "12. Technical Limitations & Future Research Perspectives")
    box13 = slide13.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf13 = box13.text_frame; tf13.word_wrap = True
    tf13.paragraphs[0].text = "Academic Transparency & Path Forward:"
    tf13.paragraphs[0].font.size = Pt(22); tf13.paragraphs[0].font.bold = True; tf13.paragraphs[0].font.color.rgb = CYAN_ACCENT

    lim_data = [
        ("DEM Spatial Resolution Constraints", "Current topography modeling is bounded by open 30-meter DEM resolution, which inevitably abstracts micro-urban terrain structures like local curbs, culverts, and individual flood defense retaining walls.", AMBER_ACCENT),
        ("Calibration & Validation Infrastructure", "Due to an absence of dense real-time municipal IoT river stream gauges, model calibration relies on historical news documentation and municipal flood reporting records.", AMBER_ACCENT),
        ("Future Scope 1: High-Resolution Drone LiDAR", "Upgrading topographic baseline rasters to sub-meter drone-derived LiDAR surface models for micro-hydrological precision.", CYAN_ACCENT),
        ("Future Scope 2: Real-Time IoT Sensor Ingestion & 2D Hydrodynamics", "Integrating real-time telemetry river gauge data streams and embedding fully coupled 2D hydrodynamic simulation engines.", CYAN_ACCENT)
    ]
    for l_head, l_body, col in lim_data:
        p_l = tf13.add_paragraph()
        p_l.text = f"•  {l_head}:"
        p_l.font.size = Pt(18); p_l.font.bold = True; p_l.font.color.rgb = col
        p_l.space_before = Pt(12)
        p_lb = tf13.add_paragraph()
        p_lb.text = f"    {l_body}"
        p_lb.font.size = Pt(16); p_lb.font.color.rgb = TEXT_WHITE
    add_speaker_notes(slide13, "In maintaining academic rigor and scientific integrity, it is essential to highlight our research boundaries. Presently, spatial precision is governed by our 30-meter Digital Elevation Model resolution, which naturally abstracts macro-topographical features while missing sub-meter drainage structures like curbs or stormwater culverts. Looking forward, our architecture is built modularly to ingest 1-meter drone-derived LiDAR datasets and live IoT water level telemetry sensors, which will eventually allow real-time coupling with complex finite-element hydrodynamic physics equations.")

    print("Synthesizing Slide 14: Conclusion & Contributions...")
    slide14 = prs.slides.add_slide(blank_layout)
    apply_slide_base(slide14, "13. Conclusion & Core Academic Contributions")
    box14 = slide14.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(5.5))
    tf14 = box14.text_frame; tf14.word_wrap = True
    tf14.paragraphs[0].text = "Synthesis of Dissertation Contributions:"
    tf14.paragraphs[0].font.size = Pt(22); tf14.paragraphs[0].font.bold = True; tf14.paragraphs[0].font.color.rgb = CYAN_ACCENT

    concl = [
        "Successfully unified traditional Multi-Criteria Decision Analysis (AHP) with modern interactive 3D WebGL Digital Twin technology.",
        "Delivered verified empirical impact exposure calculations across 180,307 buildings and 7,445.90 km of road transport networks in Pune.",
        "Established an innovative, explainable GeoAI operational blueprint designed explicitly for Smart City governance and resilient urban disaster coordination."
    ]
    for c in concl:
        p_c = tf14.add_paragraph()
        p_c.text = f"\n✓   {c}"
        p_c.font.size = Pt(18); p_c.font.bold = True; p_c.font.color.rgb = TEXT_WHITE

    p_thanks = tf14.add_paragraph()
    p_thanks.text = "\n\nTHANK YOU FOR YOUR ATTENTION.\nOPEN FOR EVALUATION & DEFENSE Q&A."
    p_thanks.font.size = Pt(24); p_thanks.font.bold = True; p_thanks.font.color.rgb = CYAN_ACCENT
    p_thanks.alignment = PP_ALIGN.CENTER
    add_speaker_notes(slide14, "In conclusion, GeoNarrative AI successfully demonstrates that standard GIS hazard evaluation can be transformed from static archival maps into an intelligent, interactive, and actionable Smart City decision support instrument. By combining a validated five-factor AHP hazard model with empirical exposure quantification across more than 180,000 structures in Pune, this dissertation provides a reproducible and robust framework for urban flood resilience. I sincerely thank the committee for your time and attention, and I now welcome any questions or suggestions.")

    output_path = "GeoNarrative_AI_MSc_Defense_Presentation.pptx"
    prs.save(output_path)
    print(f"\nSUCCESS! Fully functional presentation saved to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_deck()
